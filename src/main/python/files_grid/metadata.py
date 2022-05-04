import docx
from dateutil.parser import parse
import glob
import io
import logging
import os
import re
import shutil
import socket
import time
import traceback
import urllib
import zipfile
from datetime import datetime
import PyPDF2
import docx2txt
import googletrans
import pandas
import requests
import bs4
import textract
import usaddress
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.wait import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import pycountry
from geopy.geocoders import Nominatim
import pycld2 as cld2
import counters
import extract_pdf2
import regex_all
from pptx import Presentation
from tika import parser
from docx import *
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
from urllib.parse import unquote

class Metadata:

    # chrome_driver_path ="/home/ubuntu/kiran/chrome-headless/drivers/chromedriver"
    # chrome_driver_path="/Users/msk/Pycharm_Pro/beautiful_soup/drivers/chromedriver"
    # driver = webdriver.Chrome(service_log_path='/dev/null')
    # chromeOptions = Options()
    # chromeOptions.headless = True
    # browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions,service_log_path='/dev/null')
    # os.environ['WDM_LOG_LEVEL'] = '0'

    class resp_url:
        resp = ""
        url = ""
        domain = ""
        pcl = ""
        def __init__(self, resp, url, domain, pcl):
            self.resp = resp
            self.url = url
            self.domain = domain
            self.pcl = pcl

    def __init__(self, resume_rgx=''):
        self.country_full_name_list = []
        self.country_2c_list = []
        self.country_3c_list = []
        for ctry_obj in list(pycountry.countries):
            self.country_full_name_list.append(ctry_obj.name.lower())
            self.country_2c_list.append(ctry_obj.alpha_2)
            self.country_3c_list.append(ctry_obj.alpha_3)
        self.country_chars= ["-"," ",":",","]
        self.resume_pattern = re.compile(resume_rgx)
        # self.pages = []
        # self.easy_list = []

    def is_date(self, string, fuzzy=False):
        """
        Return whether the string can be interpreted as a date.

        :param string: str, string to check for date
        :param fuzzy: bool, ignore unknown tokens in string if True
        """
        try:
            parse(string, fuzzy=fuzzy)
            return True

        except ValueError:
            return False

    def is_resume(self, url):
        othr_fmts = ['.rtf','.csv','.xml','.json']
        flag = False
        try:
            if url.endswith('.pdf'):
               text = self.get_text_from_pdf(url)
               return self.contains_resume(text)
            if url.endswith('.docx'):
               text = self.get_text_from_docx(url)
               return self.contains_resume(text)
            if url.endswith('doc'):
                text = self.get_text_from_doc(url)
                return self.contains_resume(text)
            if url.endswith('.xlsx') or url.endswith('.xls'):
               df = self.get_df_from_excel(url)
               return self.contains_resume_in_df(df)
            if any(url.endswith(s) for s in othr_fmts):
                text = self.get_text_from_othr_file(url)
                return self.contains_resume(text)
            if url.endswith('.pptx'):
               text = self.get_text_from_pptx(url)
               return self.contains_resume(text)
            if url.endswith('.ppt'):
               text = self.get_text_from_ppt(url)
               return self.contains_resume(text)
            if url.endswith('.zip'):
                return self.does_zip_contain_resume(url)
            resp = requests.get(url, timeout=5)
            self.req_resp_size_logger(resp)
            if resp.status_code == 200:
                flag = self.contains_resume(resp.text)
        except Exception as e:
            logging.error("unable to parse url {}: {}".format(url,e))
        return flag

    def get_text_from_file_url(self, url):
        othr_fmts = ['.rtf','.csv','.xml','.json']
        try:
            if url.endswith('.pdf'):
               return self.get_text_from_pdf(url)
            if url.endswith('.docx'):
               return self.get_text_from_docx(url)
            if url.endswith('doc'):
                return self.get_text_from_doc(url)
            if url.endswith('.xlsx') or url.endswith('.xls'):
               df = self.get_df_from_excel(url)
               return self.get_text_from_df(df)
            if any(url.endswith(s) for s in othr_fmts):
                return self.get_text_from_othr_file(url)
            if url.endswith('.pptx'):
               return self.get_text_from_pptx(url)
            if url.endswith('.ppt'):
               return self.get_text_from_ppt(url)
            if url.endswith('.zip'):
                return self.get_text_from_zip_file_url(url)
        except Exception as e:
            logging.error("unable to parse url {}: {}".format(url,e))
        return ''

    def contains_resume(self, text):
        "work in progress its not 100% accurate"
        matches = re.finditer(self.resume_pattern, text)
        return any(True for _ in matches)
        # for match in matches:
        #     print(match)
        # return re.search(regex_all.resume, text) is not None

    def contains_resume_in_df(self, df):
        flag = False
        for i, j in df.iterrows():
            col_row = '{}:{}'.format(i, j)
            if None != re.search(regex_all.resume, col_row, re.IGNORECASE):
                flag = True
                break
        return flag

    def get_text_from_df(self, df):
        text_lst = []
        for i, j in df.iterrows():
            col_row = '{}:{}'.format(i, j)
            text_lst.append(col_row)
        return ' '.join(text_lst)

    def get_text_from_pdf_old(self, source):
        if self.is_url(source):
            url = source
            req = urllib.request.Request(url, headers={'User-Agent': "Magic Browser"})
            remote_file = urllib.request.urlopen(req).read()
            remote_file_bytes = io.BytesIO(remote_file)
            pdf_texts = []
            pdfreader = PyPDF2.PdfFileReader(remote_file_bytes)
            x = pdfreader.numPages
            i =0
            while i < x:
                pageobj = pdfreader.getPage(i)
                pdf_texts.append(pageobj.extractText())
                i+=1
            return ''.join(pdf_texts)
        else:
            return extract_pdf2.extract(source)

    def get_text_from_pdf(self, source):
        if self.is_url(source):
            return extract_pdf2.extract_from_url(source)
        else:
            return extract_pdf2.extract(source)

    def is_url(self, str):
        return str.startswith('http')

    def get_text_from_othr_file(self, source):
        # handles .rtf,.csv,.xml,.json files
        if self.is_url(source):
            url = source
            response = requests.get(url, timeout=5)
            self.req_resp_size_logger(response)
            with io.BytesIO(response.content) as file:
                text = file.read()
            return text
        else:
            with open(source) as file:
                text = file.read()
            return text

    def get_text_from_ppt(self, source):
        if self.is_url(source):
            url = source
            response = requests.get(url, allow_redirects=True, timeout=5)
            self.req_resp_size_logger(response)
            # with io.BytesIO(response.content) as file:
            open('temp.ppt', 'wb').write(response.content)
            parsed = parser.from_file('temp.ppt')
            return (parsed["content"])
        else:
            parsed = parser.from_file(source)
            return (parsed["content"])

    def get_text_from_pptx(self, source):
        # handles .rtf files
        if self.is_url(source):
            url = source
            pptx_text = []
            response = requests.get(url, timeout=5)
            self.req_resp_size_logger(response)
            with io.BytesIO(response.content) as file:
                prs = Presentation(file)
                # print(file)
                # print("----------------------")
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            pptx_text.append(shape.text)
            return ''.join(pptx_text)
        else:
            pptx_text = []
            prs = Presentation(source)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pptx_text.append(shape.text)
            return ''.join(pptx_text)



    def get_text_from_docx(self, source):
        data = ""
        if self.is_url(source):
            url = source
            response = requests.get(url, timeout=5)
            self.req_resp_size_logger(response)
            with io.BytesIO(response.content) as file:
                data = docx2txt.process(file)
            return data
        else:
            return docx2txt.process(source)

    # def get_text_from_doc(self, url):
    #     response = requests.get(url, allow_redirects=True, timeout=5)
    #     open('temp.doc', 'wb').write(response.content)
    #     data = textract.process('temp.doc').decode()
    #     return data

    def get_text_from_doc(self, source):
        if self.is_url(source):
            url = source
            response = requests.get(url, allow_redirects=True, timeout=5)
            self.req_resp_size_logger(response)
            # with io.BytesIO(response.content) as file:
            open('temp.doc', 'wb').write(response.content)
            data = textract.process('temp.doc').decode()
            return data
        else:
            return textract.process(source).decode()

    def get_df_from_excel(self, source):
        if self.is_url(source):
            url = source
            response = requests.get(url, timeout=5)
            self.req_resp_size_logger(response)
            data = ""
            with io.BytesIO(response.content) as file:
                data = pandas.read_excel(file)
            return data
        else:
            return pandas.read_excel(source)


    def get_address_fields(self, address):
        fields = ["Website Country","Website Zip Code","Website State","Website City","Website Street"]
        rDict = {}
        try:
            ad_list = address.split(",")
            if len(ad_list) > 2:
                i = 0
                for ad_fld in reversed(ad_list):
                    rDict[fields[i]] = ad_fld
                    i = i +1
            elif len(ad_list) == 2:
                i = 0
                for ad_fld in reversed(ad_list):
                    rDict[fields[i]] = ad_fld
                    i = i + 2
            else:
                i = 2
                for ad_fld in ad_list:
                    rDict[fields[i]] = ad_fld
        except:
            logging.error("unable to parse address: {}".format(address))
        return rDict



    def get_ip_address(self, domain):
        try:
            return socket.gethostbyname(domain)
        except:
            return ''

    def is_active(self, url):
        try:
            resp = requests.head(url, timeout=10)
            if resp.status_code in (301, 302):
                return self.is_active(resp.headers.get('Location'))
        except:
            return False
        if 200 == resp.status_code:
            return True
        return False

    def get_last_modified(self, resp, url, domain):
        lm = ''
        try:
            if resp.status_code in (301, 302):
                lm = self.get_last_modified(resp.headers.get('Location'))
            if '' == lm:
                header = resp.headers
                if 'Last-Modified' in header:
                    lm = header['Last-Modified']
            if '' == lm:
                lm = self.get_last_modied_from_google(domain)
        except:
            logging.error("no last modified date on page")
        return lm

    def get_meta_tags(self, resp):
        meta_list = []
        try:
            if resp.status_code in (301, 302):
                return self.get_meta_tags(resp.headers.get('Location'))
            lxml = bs4.BeautifulSoup(resp.text, "lxml")
            title = lxml.title.string
            # print('TITLE IS :', title)
            meta_list.append('TITLE :' + title)
            meta = lxml.find_all('meta')
            for tag in meta:
                if 'name' in tag.attrs.keys() and tag.attrs['name'].strip().lower() in ['description', 'keywords']:
                    # print ('NAME    :',tag.attrs['name'].lower())
                    # print('CONTENT :', tag.attrs['content'])
                    meta_list.append('CONTENT :'+ tag.attrs['content'])
        except:
            pass
        return '\n'.join(meta_list)

    def get_contact_info_dict(self, info):
        ret_dict = {}
        infos = info.split("\n")
        owner, email, phone, address = '', '', '', ''
        for c_info in infos:
            c_info_split = c_info.split(":")
            if "Organization" == c_info_split[0]:
                owner =  c_info_split[1]
            if "Email" == c_info_split[0]:
                email = c_info_split[1]
            if "Phone" == c_info_split[0]:
                phone = c_info.split("Phone:")[1]
            if "Mailing Address" == c_info_split[0]:
                address = c_info_split[1]

        ret_dict['Owner'] = owner
        ret_dict['E-mail'] = email
        ret_dict['Phone'] = phone
        ret_dict['Address'] = address

        return ret_dict

    def get_zip_code_from_address(self, add_line):
        zip_code_str = ''
        zip_code = ''
        if add_line is None : return zip_code
        _len = len(add_line)
        for i in range(0,_len,1):
            if i+5 <=_len and add_line[i:i+5].isdigit():
                zip_code_str= add_line[i:]
                break
        if '' != zip_code_str:
            end_idx = 0
            for i in range(4,len(zip_code_str),1):
                if not zip_code_str[i].isdigit():
                    end_idx = i
                    break
            zip_code = zip_code_str[:end_idx]
        return zip_code

    def get_county(self, zip_code):
        # get from https://www.getzips.com/zip.htm
        if '' == zip_code:
            return ''
        zip_url = "https://www.getzips.com/zip.htm"
        county = ''
        try:
            chromeOptions = Options()
            chromeOptions.headless = True
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                            options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
            browser.get(zip_url)
            browser.find_element_by_id("fldZIPCode").send_keys(zip_code)
            browser.find_element_by_xpath("//p[@align='LEFT'] //input[@name='Submit']").click()
            wait = WebDriverWait(browser, 5)
            element = wait.until(
                EC.presence_of_element_located((By.XPATH, "//table[@width='540']//tr[2]/td[3]")))

            county = element.text
        except:
            logging.info("county not found for zip-code:"+zip_code)
        return county

    def get_zip_details(self, zip_code):
        # get from https://www.getzips.com/zip.htm
        zip_url = "https://www.getzips.com/zip.htm"
        zip_dict = {}
        county = ''
        city = ''
        try:
            chromeOptions = Options()
            chromeOptions.headless = True
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                            options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
            browser.get(zip_url)
            browser.find_element_by_id("fldZIPCode").send_keys(zip_code)
            browser.find_element_by_xpath("//p[@align='LEFT'] //input[@name='Submit']").click()
            wait = WebDriverWait(browser, 5)
            element = wait.until(
                EC.presence_of_element_located((By.XPATH, "//table[@width='540']//tr[2]/td[3]")))
            county = element.text
            city = browser.find_element_by_xpath("//table[@width='540']//tr[2]/td[2]").text

        except:
            logging.info("county not found for zip-code:"+zip_code)

        zip_dict["county"] = county
        zip_dict["city"] = city

        return zip_dict


    def get_contact_info_icann(self, domain):
        contact_info = ''
        try:
            icann_url = "https://lookup.icann.org/"
            chromeOptions = Options()
            chromeOptions.headless = True
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                            options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
            browser.get(icann_url)

            try:
                element = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.ID, "input-domain")))
                element.send_keys(domain)
                element.submit()
                elements = WebDriverWait(browser, 10).until(
                    EC.presence_of_all_elements_located((By.XPATH, "//div[@class='information-panel ng-star-inserted']")))
                text = elements[1].text
                if "Contact Information" in text:
                    contact_info = elements[1].text
            finally:
                browser.quit()
        except:
            logging.info("icann search not working for domain :{}".format(domain))
        return self.get_contact_info_dict(contact_info)

    def get_about_links(self, resp, domain):
        html_string = resp.text.lower()
        about_idx = html_string.find("about")
        links = re.findall("href=[\"\'](.*?)[\"\']", html_string[about_idx-200: about_idx + 200])
        if not links:
            about_us_idx = html_string.find("about us")
            links = re.findall("href=[\"\'](.*?)[\"\']", html_string[about_us_idx-200: about_us_idx + 200])
        if not links:
            domain = "about "+ domain.replace("."," ")
            about_the_idx = html_string.find(domain)
            links = re.findall("href=[\"\'](.*?)[\"\']", html_string[about_the_idx - 200: about_the_idx + 200])
        for link in links:
            if "about" in link and link.startswith("http") : return link
        return ''

    def get_about_data(self, resp, domain):
        about_strings = ["about","about_us", "about "+ domain.replace("."," ")]
        return self.get_data_around_words(resp, about_strings, 200)



    def is_secure_helper(self, domain, protocol):
        try:
            resp = requests.head(protocol+"{}".format(domain))
            if 200 == resp.status_code:
                return "http"
            elif resp.status_code in (301, 302):
                return resp.headers.get('Location').split(":")[0]
            else :
                return ''
        except:
            return ''

    def is_secure(self, domain):
        ptk_list = ["http://", "http://www.", "https://", "https://www."]
        for ptk in ptk_list:
            ptkl = self.is_secure_helper(domain, ptk)
            if ptkl != '': return ptkl
        return 'Unknown'

    def get_resp_url_helper1(self, domain, protocol):
        try:
            url = protocol+"{}".format(domain)
            resp = requests.get(url, timeout=5)
            self.req_resp_size_logger(resp)
            if 200 == resp.status_code:
                return self.resp_url(resp, url, domain, protocol)
            elif resp.status_code in (301, 302):
                return self.resp_url(resp, resp.headers.get('Location'), domain, protocol)
            else:
                return self.resp_url('', '', '', '')
        except:
            return self.resp_url('', '', '', '')

    def get_resp_url_helper0(self, domain):
        ptk_list = ["http://", "http://www.", "https://", "https://www."]
        for ptk in ptk_list:
            resp_url = self.get_resp_url_helper1(domain,ptk)
            if resp_url.resp != '': return resp_url
        return self.resp_url('','','','')

    def get_resp_url(self, domain):
        resp_url = self.get_resp_url_helper0(domain)
        if resp_url.resp == '':
            domain_types = ['.com', '.org', '.net', '.int', '.edu', '.gov', '.mil']
            d_suf = domain[-4:]
            if d_suf not in domain_types:
                for d_type in domain_types:
                    resp_url = self.get_resp_url_helper0('{}{}'.format(domain, d_type))
                    if resp_url.resp != '': return resp_url
        return resp_url


    def get_language(self, resp):
        lang = None
        try:
            text = resp.text.lower()
            soup = bs4.BeautifulSoup(text, 'lxml')
            if 'lang' in soup.html.attrs:
                lang = googletrans.LANGUAGES.get(soup.html.attrs['lang'])
            else:
                words = soup.text.split(" ")
                twenty_words = ' '.join(words[:20])
                translator = googletrans.Translator()
                lang = googletrans.LANGUAGES.get(translator.detect(twenty_words).lang)
            if None == lang:
                lang = soup.html.attrs['lang']
        except:
            lang = 'Unknown'
        return lang

    def get_languages(self, resp):
        lgs = set()
        try:
            text = resp.text.lower()
            soup = bs4.BeautifulSoup(text, 'lxml')
            if 'lang' in soup.html.attrs:
                lgs.add(soup.html.attrs['lang'])
            else:
                words = soup.text.split(" ")
                twenty_words = ' '.join(words[:20])
                lgs.update(self.get_languages_from_text(twenty_words))
        except:
            return ','.join(lgs)

    def get_code_languages(self, resp_url):
        pgm_dict = {".js": "JavaScript",
                    ".php": "PHP: Hypertext Preprocessor",
                    ".css": "Cascading Style Sheets",
                    ".jsp": "Jakarta Server Pages",
                    ".py": "Python",
                    ".sql": "Structured Query Language",
                    ".asp": "Active Server Pages",
                    ".net": ".NET Framework"}
        c_lags = set()
        c_lags.add("HTML")
        text = self.get_page_source(resp_url)
        pgm_lgs_rgx ="\.js('|\")|\.php('|\")|\.css('|\")|\.jsp('|\")|\.py('|\")|\.sql('|\")|\.asp('|\")|\.net('|\")"
        match_itr = re.finditer(pgm_lgs_rgx, text)
        for itr in match_itr:
            c_lags.add(pgm_dict.get(itr.group()[:-1]))
        return ', '.join(c_lags)


    def get_page_source(self, resp_url):
        page = ''
        try:
            chromeOptions = Options()
            chromeOptions.headless = True
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                            options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
            browser.get(resp_url.url)
            page = browser.page_source
            browser.quit()
        except:
            logging.error("unable to get page source for:"+resp_url.url)
        return page


    def get_languages_from_text(self, text):
        lgs = set()
        detected_lgs = cld2.detect(text, returnVectors=True)
        for obj in detected_lgs:
            if type(obj) is tuple:
                for lg_obj in obj:
                    if type(lg_obj) is tuple:
                        for lg in lg_obj:
                            if type(lg) is str and len(lg) > 2 and 'Unknown' != lg: lgs.add(lg)
        return lgs


    def is_payments_available(self, resp):
        html_string = resp.text.lower()
        return html_string.find("payment") != -1 or html_string.find("pay ") != -1


    def get_ip_country(self, domain):
        # logger = logging.getLogger('selenium.webdriver.remote.remote_connection')
        # logger.setLevel(logging.WARNING)
        county = ''
        google_url = "https://www.site24x7.com/find-website-location.html"

        chromeOptions = Options()
        chromeOptions.headless = True
        # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
        #                            options=chromeOptions)
        # browser = webdriver.Chrome(service_log_path='NULL')
        browser = webdriver.Chrome(ChromeDriverManager().install(),options=chromeOptions)
        try:
            browser.get(google_url)

            browser.find_element_by_id("hostName").clear()
            browser.find_element_by_id("hostName").send_keys(domain)
            browser.find_element_by_id("cmdbtn").click()
            time.sleep(5)

            wait = WebDriverWait(browser, 10)
            element = wait.until(
                EC.presence_of_element_located((By.XPATH, "//div[@id='fpRowTempCN2']/div[2]")))

            county = element.text
        except:
            logging.info("country not found for domain: " + domain)
        finally:
            browser.quit()
        return county


    def is_captcha_enabled(self, resp):
        html_string = resp.text.lower()
        return html_string.find("captcha") != -1

    def get_links_around_words(self, resp, words):
        html_string = resp.text.lower()
        _links =[]
        # logging.info(html_string)
        for word in words:
            word_idx = html_string.find(word)
            if -1 == word_idx: continue
            # logging.info(html_string[word_idx - 200: word_idx + 200])
            links = re.findall("href=[\"\'](.*?)[\"\']", html_string[word_idx - 200: word_idx + 200])
            if not not links:
                # logging.info(links)
                _links.extend(links)
        return _links

    def get_data_around_words(self, resp, words, data_len):
        html_string = resp.text.lower()
        # logging.info(html_string)
        for word in words:
            word_idx = html_string.find(word)
            if -1 == word_idx: continue
            # logging.info(html_string[word_idx - 200: word_idx + 200])
            return html_string[word_idx - data_len: word_idx + data_len]


    def is_signup_login_available(self, resp):
        search_words = ["sign in", "signin", "sign up", "signup", "loginin", "log in"]
        return not not self.get_links_around_words(resp, search_words)

    def get_offering(self, resp):
        html_string = resp.text.lower()
        soup = bs4.BeautifulSoup(html_string,'lxml')
        offerings = ["products", "services"]
        ret_list = []
        for offering in offerings:
            if soup.text.find(offering) != -1: ret_list.append(offering)
        # default to services
        if not ret_list:
            ret_list.append(offerings[1])
        return ','.join(ret_list)



    def get_category(self, domain):
        out_list = []
        try:
            icann_url = "https://www.safedns.com/check/"
            chromeOptions = Options()
            chromeOptions.headless = True
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                            options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
            browser.get(icann_url)

            try:
                element = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.XPATH, "//input[@class='search-block__input']")))
                element.send_keys(domain+".com")
                element = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.ID, "check_data")))
                element.click()
                elements = WebDriverWait(browser, 10).until(
                    EC.presence_of_all_elements_located((By.CLASS_NAME, "check-page__stg-two-categories-item")))
                for elem in elements:
                    out_list.append(elem.text)
            finally:
                browser.quit()
        except:
            # traceback.logging.info_exc()
            logging.info("safedns search not working for domain :{}".format(domain))
        return ','.join(out_list)

    def search_google(self, text):
        # below method should be me results from google search
        result = ''
        try:
            google_url = "https://www.google.com/"
            chromeOptions = Options()
            chromeOptions.headless = True
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                            options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
            browser.get(google_url)

            try:
                element = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.XPATH, "//input[@title='Search']")))
                element.send_keys(text)
                element = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.XPATH, "//input[@title='Search']")))
                element.send_keys(Keys.ENTER)
                element = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.XPATH, "//div[@id='rcnt']")))
                result = element.text
            finally:
                browser.quit()
        except:
            traceback.print_exc()
            logging.info("google search not working for text :{}".format(text))
        return result

    def search_google_bs4(self, text):
        # below method should be me results from google search
        result = ''
        try:
            google_url = 'https://google.com/'
            res = requests.get(google_url+'search?q='+text, timeout = 5)
            res.raise_for_status()
            soup = bs4.BeautifulSoup(res.text, 'lxml')
            result = soup.text
        except:
            traceback.print_exc()
            logging.info("google search not working for text :{}".format(text))
        return result

    def get_age_restriction_gsearch(self, domain):
        results = self.search_google_bs4("what is the age restriction of "+domain).split("\n")
        for result in results:
            _len = len(result)
            for i in range(0,_len,1):
                if i+3 <= _len and not result[i].isdigit() and result[i+1].isdigit() and result[i+2].isdigit() and not result[i+3].isdigit():
                    return result[i+1:i+3]+"+"
        return ''


    def get_whatsapp_number(self, resp):
        w_num = []
        whatsapp_regex = "\+[0-9-\s]+"
        soup = bs4.BeautifulSoup(resp.text,'lxml')
        lower_text = soup.text.lower()
        if -1 != lower_text.find("whatsapp"):
            w_num = re.findall(whatsapp_regex,lower_text)
        return  ','.join(w_num)

    def get_nature_of_content(self, resp):
        noc_types = {"Vedios":"video", "Images":"img", "Links":"href"}
        noc_list = []
        for key,val in noc_types.items():
            if resp.text.lower().find(val) != -1:
                noc_list.append(key)
        noc_list.append("Text")
        return ','.join(noc_list)

    def get_file_types(self, resp):
        file_types = [".pdf", ".xml", ".ppt", ".xls", ".doc", ".docx", ".rtf"]
        return self.get_links_around_words(resp, file_types)

    def get_google_search_result_links(self, text):
        ret_list = []
        # below method should be me results from google search
        try:

            google_url = "https://www.google.com/"
            chromeOptions = Options()
            chromeOptions.headless = True
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                            options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)

            browser.implicitly_wait(5)
            browser.get(google_url)


            try:
                element = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.XPATH, "//input[@title='Search']")))
                element.send_keys(text)
                element = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.XPATH, "//input[@title='Search']")))
                element.send_keys(Keys.ENTER)
                list_browser = WebDriverWait(browser, 10).until(
                    EC.presence_of_element_located((By.ID, "rso")))

                lia = list_browser.find_elements_by_tag_name("a")

                for i in range(0, len(lia), 1):
                    if "google" not in list_browser.find_elements_by_tag_name("a")[i].get_attribute('href'):
                        # print(list_browser.find_elements_by_tag_name("a")._getitem_(i).get_attribute("href"))
                        ret_list.append(list_browser.find_elements_by_tag_name("a")[i].get_attribute('href'))

            finally:
                browser.quit()

        except:
            logging.info("google search not working for text :{}".format(text))
        return ret_list

    def get_primary_address(self, domain):
        p_address = ''
        result =self.search_google_bs4("what is the address of "+domain)
        # print(result)
        for line in result.split("\n"):
            if '' != line.strip() and self.parse_address(line):
                p_address= line
                break
        p_address = self.format_address(p_address)
        if '' == p_address:
            links = self.get_google_search_result_links("what is the address of " + domain)
            out_link = ''
            for link in links:
                if "contact" in link.lower():
                    out_link = link
                    break
            if '' == out_link:
                for link in links:
                    if "about" in link.lower():
                        out_link = link
                        break
            # print(out_link)
            if '' != out_link:
                p_address = self.get_address_from_page(out_link)
        return p_address

    def parse_address(self, line):
        lower_line = line.lower()
        tot_line_idx = len(lower_line)
        if "address:" in lower_line: return True
        if "ip address" in lower_line: return False
        if self.hasNumbers(lower_line):
            for c_fn in self.country_full_name_list:
                if c_fn in lower_line:
                    return True
            for c_n in self.country_2c_list:
                idx = line.find(c_n)
                if -1 != idx and line[idx - 1] in self.country_chars and (
                        idx+2 == tot_line_idx or line[idx + 2] in self.country_chars):
                    return True
            for c_n in self.country_3c_list:
                idx = line.find(c_n)
                if -1 != idx and line[idx - 1] in self.country_chars and (
                        idx+3 == tot_line_idx or line[idx + 3] in self.country_chars):
                    return True
        return False

    def hasNumbers(self, inputString):
        return any(char.isdigit() for char in inputString)

    def get_address_from_page(self, link):
        add_list = []
        address_pattern = "\d{1,5}\s(\\b\w*\\b\s)\w*(\.|,)"
        html =''
        try:
            html = requests.get(link, timeout=5).text
        except:
            html = ''
        soup = bs4.BeautifulSoup(html,'lxml')
        add_itr = re.finditer(address_pattern,soup.text)
        zip_code_pattern = "\d{5}([-]|\s*)?(\d{4})?"
        for m in add_itr:
            s_idx = m.start()
            nl_idx = soup.text[s_idx:].find("\n")
            add_line = soup.text[s_idx:s_idx+nl_idx]
            zip_code_idx = None
            for z in re.finditer(zip_code_pattern, add_line):
                zip_code_idx = z.end()
                break
            if zip_code_idx is not None:
                add_list.append(add_line[:zip_code_idx])
        return add_list[0] if not not add_list else ''

    def get_geolocation(self, address):
        geolocator = Nominatim(user_agent="myGeocoder")
        g_loc = ''
        if address.strip() == '': return g_loc
        try:
            location = geolocator.geocode(address)  # Searching for place
            latitude = location.latitude  # Getting Latitude
            longitude = location.longitude  # Getting Longitude
            g_loc = 'latitude: {}, longitude: {}'.format(latitude, longitude)
        except Exception as e:
            logging.error("geolocation not found for address:"+address)
            # traceback.print_exc()
        return g_loc

    def get_primary_address_dict_old(self, line, add_dict):

        geolocator = Nominatim(user_agent="myGeocoder")
        match = re.search(regex_all.zip_code, line)
        if match is not None:
            geo_details = geolocator.geocode(match.group())
            if geo_details is not None:
                add_dict['Primary Address'] = geo_details.address
                add_dict['Primary Address Geolocation'] = 'latitude: {}, longitude: {}'.format(geo_details.latitude,
                                                                                               geo_details.longitude)
                add_dict['Primary Address County'] = geo_details.raw['address']['county']
                add_dict['Primary Address City'] = geo_details.raw['address']['state']
                add_dict['Primary Address Country'] = geo_details.raw['address']['country']
                add_dict['Primary Address Zip Code'] = geo_details.raw['address']['postcode']

        return add_dict

    # geocode address format
    # House Number, Street Direction, Street Name, Street Suffix, City, State, Zip, Country
    # address is a String e.g. 'Berlin, Germany'
    # addressdetails=True does the magic and gives you also the details
    def get_primary_address_dict(self, line, add_dict):
        _match = re.search(regex_all.us_address, line)
        if _match is None : return add_dict
        src_address = _match.group()
        usa_address = usaddress.parse(src_address)
        print(usa_address)
        geo_add = []
        if usa_address:
            [geo_add.append(add_tup[0].replace(',', '')) for add_tup in usa_address]
            geo_add.append('USA')
            print(geo_add[-4:])
        geo_address = ', '.join(geo_add[-4:])
        geolocator = Nominatim(user_agent="myGeocoder")
        geo_details = geolocator.geocode(geo_address, addressdetails=True)
        if geo_details is not None:
            add_dict['Primary Address'] = geo_details.address
            add_dict['Primary Address Geolocation'] = 'latitude: {}, longitude: {}'.format(geo_details.latitude, geo_details.longitude)
            add_dict['Primary Address County'] = geo_details.raw['address']['county']
            add_dict['Primary Address City'] = geo_details.raw['address']['state']
            add_dict['Primary Address Country'] = geo_details.raw['address']['country']
            add_dict['Primary Address Zip Code'] = geo_details.raw['address']['postcode']

        return add_dict

    def do_ads_exits(self, resp):
        return "TRUE" if -1 != resp.text.find("ads") else "FALSE"

    def get_web_code_lang_version(self, resp_url):
        soup = bs4.BeautifulSoup(resp_url.resp.text,'lxml')
        doc_type =self.doctype(soup).lower()
        if -1 != doc_type.find("html 4") or -1 != doc_type.find("xhtml 1"): return "HTML 4"
        if -1 != doc_type.find("html 3"): return "HTML 3"
        if -1 != doc_type.find("html 2"): return "HTML 2"
        if -1 != doc_type.find("html 1"): return "HTML 1"
        return "HTML 5"

    def doctype(self ,soup):
        items = [item for item in soup.contents if isinstance(item, bs4.Doctype)]
        return items[0] if items else None

    def get_business_hours(self, domain):
        for match in re.finditer("\.", domain):
            if match.start() is not None:
                i = match.start()
        domain = domain[:i]
        # not_allowed_cat =["Computers & Internet"]
        # if ctgy in not_allowed_cat: return ''
        google_url = "https://www.google.com/maps/@12.9333774,77.5963265,15z"

        chromeOptions = Options()
        chromeOptions.headless = True
        # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
        #                            options=chromeOptions)
        browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
        b_hours = ''
        try:

            browser.get(google_url)
            browser.find_element_by_id("searchboxinput").send_keys(domain)
            browser.find_element_by_id("searchbox-searchbutton").click()

            try:
                wait1 = WebDriverWait(browser, 15)
                element1 = wait1.until(
                    EC.visibility_of_element_located((By.XPATH, "//a[@class='a4gq8e-aVTXAb-haAclf-jRmmHf-hSRGPd']")))
                element1.click()
            except:
                wait2 = WebDriverWait(browser, 15)
                element2 = wait2.until(
                    EC.visibility_of_element_located((By.XPATH, "//span[@aria-label='Show open hours for the week']")))
                element2.click()

                b_hours = browser.find_element_by_xpath("//table[@class='y0skZc-jyrRxf-Tydcue NVpwyf-qJTHM-ibL1re']").text
                browser.quit()
                return b_hours

            wait2 = WebDriverWait(browser, 15)
            element2 = wait2.until(
                EC.visibility_of_element_located((By.XPATH, "//span[@aria-label='Show open hours for the week']")))
            element2.click()

            b_hours = browser.find_element_by_xpath("//table[@class='y0skZc-jyrRxf-Tydcue NVpwyf-qJTHM-ibL1re']").text
            browser.quit()
        except:
            logging.error("unable to get business hours for domain {}".format(domain))
            # traceback.print_exc()
        return b_hours

    def is_mobile_friendly(self, domain):
        google_url = "https://technicalseo.com/tools/mobile-friendly/"
        chromeOptions = Options()
        chromeOptions.headless = True
        ret_value = 'FALSE'
        try:
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                            options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
            browser.get(google_url)
            time.sleep(7)
            element = WebDriverWait(browser, 10).until(
                EC.presence_of_element_located((By.ID, "onetrust-accept-btn-handler")))
            element.click()
            element = WebDriverWait(browser, 10).until(
                EC.presence_of_element_located((By.XPATH, "//textarea[@id='input_0']"))
            )
            element.send_keys(domain)

            element = WebDriverWait(browser, 10).until(
                EC.presence_of_element_located((By.XPATH, "//body/div[2]/div[1]/section[1]/div[1]/div[1]/div[1]/form[1]/div[1]/div[2]/md-input-container[1]/button[1]"))
            )
            element.click()
            element = WebDriverWait(browser, 50).until(
                EC.presence_of_element_located((By.XPATH, "/html/body/div[2]/div[1]/section/div/div[2]/div/md-table-container/table/tbody/tr/td[2]/span")))

            text = element.find_element_by_xpath(
                "/html/body/div[2]/div[1]/section/div/div[2]/div/md-table-container/table/tbody/tr/td[2]/span").get_attribute(
                'textContent')
            browser.quit()
            if "yes" == text: ret_value = "TRUE"
        except:
            logging.error("error during mobile friendly check: "+domain)
        return ret_value


    def format_address(self, address_line):
        add_line =''
        for _match in re.finditer(regex_all.address_line,address_line):
            add_line = address_line[_match.start():_match.end()]
            break

        # zip_match = re.search("[0-9]\.",add_line)
        # zip_match = re.search(regex_all.zip_code,add_line)
        # if zip_match is not None:
        #     add_line = add_line[:zip_match.end()]
        country_match = re.finditer(regex_all.capital_country, add_line)
        idx = len(add_line)
        for itr in country_match:
            idx = itr.end()

        if idx != len(add_line):
            cnty_match = re.search("[a-z]\.", add_line)
            if cnty_match is not None:
                add_line = add_line[:cnty_match.end()]
                
        return add_line[:idx]

    def get_primary_address_fields_gsearch(self, domain):
        fields = ["Primary Address", "Primary Address Geolocation", "Primary Address County", "Primary Address City",
                  "Primary Address Country", "Primary Address Zip Code"]
        r_dict = {fields[0]: self.get_primary_address(domain)}
        r_dict[fields[1]] = self.get_geolocation(r_dict[fields[0]])
        zip_code = ''
        zip_dict = {}
        zip_match = re.search(regex_all.zip_code, r_dict[fields[0]])
        if zip_match is not None:
            zip_code = r_dict[fields[0]][zip_match.start():zip_match.end()]
            zip_dict = self.get_zip_details(zip_code)

        r_dict[fields[2]] = zip_dict.get("county")
        r_dict[fields[3]] = zip_dict.get("city")
        r_dict[fields[4]] = self.get_country_from_address_line(r_dict[fields[0]], zip_dict.get("city"))
        r_dict[fields[5]] = zip_code

        return r_dict


    def get_country_from_address_line(self, line, city):
        lower_line = line.lower()
        for c_n in self.country_full_name_list:
            idx = lower_line.find(c_n)
            if -1 != idx:
                return line[idx:idx + len(c_n)]
        for c_n in self.country_2c_list:
            idx = line.find(c_n)
            if -1 != idx and city not in line:
                return line[idx:idx+len(c_n)]
        for c_n in self.country_3c_list:
            idx = line.find(c_n)
            if -1 != idx:
                return line[idx:idx + len(c_n)]
        return ''

    def get_primary_email(self, text):
        email = ''
        try:
            m = re.search(regex_all.email, text)
            if m: email = m.group(0)
        except:
            return email
        return email

    def get_primary_phone(self, text):
        phone = ''
        try:
            m = re.search(regex_all.phone, text)
            if m: phone = m.group(0)
        except:
            return phone
        return phone

    def get_google_rating(self, d):
        rating = ''
        try:
            google_url = "https://www.google.com/maps/@12.9333774,77.5963265,15z"
            chromeOptions = Options()
            chromeOptions.headless = True
            # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
            #                        options=chromeOptions)
            browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
            browser.get(google_url)
            browser.find_element_by_id("searchboxinput").send_keys(d)
            browser.find_element_by_id("searchbox-searchbutton").click()

            try:
                element = WebDriverWait(browser, 10).until(
                EC.presence_of_element_located((By.XPATH, "//a[@class='a4gq8e-aVTXAb-haAclf-jRmmHf-hSRGPd']")))
                element.click()
            except:
                element2 = WebDriverWait(browser, 10).until(
                EC.visibility_of_element_located((By.XPATH, "//div[@jsaction='pane.rating.moreReviews']  //span[@class='aMPvhf-fI6EEc-KVuj8d']")))
                rating = "{} out of 5".format(element2.text)
                browser.quit()
                return rating   

            element2 = WebDriverWait(browser, 10).until(
            EC.visibility_of_element_located((By.XPATH, "//div[@jsaction='pane.rating.moreReviews']  //span[@class='aMPvhf-fI6EEc-KVuj8d']")))
            rating = "{} out of 5".format(element2.text)
            browser.quit()
        except:
            logging.error("google rating not found for domain: "+d)
        return rating

    def is_hiring(self, resp):
        words = ["careers", "jobs", "hiring", "postings", "placements", "employment"]
        result = self.get_data_around_words(resp, words, 100)
        return "FALSE" if result is None else "TRUE"


    def get_last_modied_from_google(self, domain):
        last_modified = ''
        google_url = "https://www.google.com.ua/search?q=site%3Awww."+domain+"&biw=1855&bih=" \
                     "916&source=lnt&tbs=cdr%3A1%2Ccd_min%3A1%2F1%2F2000%2Ccd_max%3A&tbm="

        chromeOptions = Options()
        chromeOptions.headless = True
        # browser = webdriver.Chrome(executable_path=os.getenv("chrome_driver_path"),
        #                            options=chromeOptions)
        browser = webdriver.Chrome(ChromeDriverManager().install(), options=chromeOptions)
        browser.get(google_url)
        list = browser.find_elements_by_css_selector(".MUxGbd.wuQ4Ob.WZ8Tjf")
        dates = []
        for obj in list:
            dates.append(obj.text.replace(" —",""))

        if dates:
            dates.sort(key=lambda date: datetime.strptime(date, '%d-%b-%Y'))
            last_modified = dates[len(dates) -1]

        return last_modified

    def does_zip_contain_resume(self, url):
        resume_extns= ['rtf','csv','xml','json','pdf','docx','doc','xlsx','xls','pptx','ppt','zip']
        zip_file_name ='/tmp/zip_file.zip'
        extract_path = '/tmp/zip_file'
        zip_dir_files = '/tmp/zip_file/*'
        if not os.path.exists(extract_path):
            os.mkdir(extract_path)
        else:
            for f in glob.glob(zip_dir_files):
                os.remove(f)
        # wget.download(url, zip_file_name)
        resp = requests.get(url, allow_redirects=True, timeout=5)
        self.req_resp_size_logger(resp)
        open(zip_file_name, 'wb').write(resp.content)
        self.extract_zip_files(zip_file_name, extract_path)
        # with zipfile.ZipFile(zip_file_name, 'r') as zip_ref:
        #     zip_ref.extractall(extract_path)
        for file in glob.glob(zip_dir_files):
            file_tkns = file.split('.')
            extn = file_tkns[len(file_tkns) -1]
            if extn in resume_extns:
                flag = self.is_file_resume(file, extn)
                if flag: return flag


    def extract_zip_files(self, zip_file_name, target_dir):
        with zipfile.ZipFile(zip_file_name) as zip_file:
            for member in zip_file.namelist():
                filename = os.path.basename(member)
                # skip directories
                if not filename:
                    continue
                # copy file (taken from zipfile's extract)
                source = zip_file.open(member)
                target = open(os.path.join(target_dir, filename), "wb")
                with source, target:
                    shutil.copyfileobj(source, target)
        for f in glob.glob(zip_file_name):
            os.remove(f)

    def get_text_from_zip_file_url(self, url):
        texts =[]
        resume_extns= ['rtf','csv','xml','json','pdf','docx','doc','xlsx','xls','pptx','ppt','zip']
        zip_file_name ='/tmp/zip_file.zip'
        extract_path = '/tmp/zip_file'
        zip_dir_files = '/tmp/zip_file/*'
        if not os.path.exists(extract_path):
            os.mkdir(extract_path)
        else:
            for f in glob.glob(zip_dir_files):
                os.remove(f)
        # wget.download(url, zip_file_name)
        resp = requests.get(url, allow_redirects=True, timeout=5)
        self.req_resp_size_logger(resp)
        open(zip_file_name, 'wb').write(resp.content)
        self.extract_zip_files(zip_file_name, extract_path)
        # with zipfile.ZipFile(zip_file_name, 'r') as zip_ref:
        #     zip_ref.extractall(extract_path)
        for file in glob.glob(zip_dir_files):
            file_tkns = file.split('.')
            extn = file_tkns[len(file_tkns) -1]
            if extn in resume_extns:
                texts.append(self.get_text_from_file(file, extn))
        return '\n'.join(texts)

    def get_text_from_file(self, file, extn):
        othr_fmts = ['rtf', 'csv', 'xml', 'json']
        try:
            if extn == 'pdf':
                return self.get_text_from_pdf(file)
            if extn == 'docx':
                return self.get_text_from_docx(file)
            if extn == 'doc':
                return self.get_text_from_doc(file)
            if extn in ['xlsx','xls']:
                df = self.get_df_from_excel(file)
                return self.get_text_from_df(df)
            if any(extn == (s) for s in othr_fmts):
                return self.get_text_from_othr_file(file)
            if extn == 'pptx':
                return self.get_text_from_pptx(file)
            if extn == 'ppt':
                return self.get_text_from_ppt(file)
        #     zip file within zip file not handled
        except Exception as e:
            logging.error("unable to parse extn {}: {}".format(extn, e))
        return ''


    def is_file_resume(self, file, extn):
        othr_fmts = ['rtf', 'csv', 'xml', 'json']
        flag = False
        try:
            if extn == 'pdf':
                text = self.get_text_from_pdf(file)
                # print('pdf text {}'.format(len(text)))
                return self.contains_resume(text)
            if extn == 'docx':
                text = self.get_text_from_docx(file)
                # print('docx text {}'.format(len(text)))
                return self.contains_resume(text)
            if extn == 'doc':
                text = self.get_text_from_doc(file)
                # print('doc text {}'.format(len(text)))
                return self.contains_resume(text)
            if extn in ['xlsx','xls']:
                df = self.get_df_from_excel(file)
                # print('excel text {}'.format(len(df)))
                return self.contains_resume_in_df(df)
            if any(extn == (s) for s in othr_fmts):
                text = self.get_text_from_othr_file(file)
                # print('other text {}'.format(len(text)))
                return self.contains_resume(text)
            if extn == 'pptx':
                text = self.get_text_from_pptx(file)
                # print('pptx text {}'.format(len(text)))
                return self.contains_resume(text)
            if extn == 'ppt':
                text = self.get_text_from_ppt(file)
                # print('ppt text {}'.format(len(text)))
                return self.contains_resume(text)
            if extn == 'zip':
                return self.does_zip_contain_resume(file)
        except Exception as e:
            logging.error("unable to parse extn {}: {}".format(extn, e))
        return flag

    def req_resp_size_logger(self, response):
        try:
            method_len = len(response.request.method)
            url_len = len(response.request.url)
            headers_len = len('\r\n'.join('{}{}'.format(k, v) for k, v in response.request.headers.items()))
            body_len = len(response.request.body if response.request.body else [])
            req_size = method_len + url_len + headers_len + body_len
            resp_size = len(response.content)
            # logging.info(f'Request size in bytes: {req_size}')
            # logging.info(f'Response size in bytes: {resp_size}')
            counters.set_req_size(req_size)
            counters.set_resp_size(resp_size)
        except Exception as e:
            logging.error("unable to log size of req/resp {}".format(e))

    def get_resp_from_url(self, url):
        try:
            resp = requests.get(url, timeout = 5)
            self.req_resp_size_logger(resp)
            if 200 == resp.status_code:
                return resp
        except Exception as e:
            logging.error("url not reachable {}: {}".format(url,e))

    def get_file_name(self, url):
        file_name = ''
        tokens = url.split("//")
        if len(tokens) > 0:
            names = tokens[1].split('/')
            name_len = len(names)
            if name_len > 0: 
                file_name = names[name_len -1]
        return unquote(file_name)

    def get_file_description(self, url):
        return 'no metadata found'

    def get_file_extension(self, url):
        extn = ''
        tokens = url.split('.')
        tok_len = len(tokens)
        if tok_len > 0:
            extn = tokens[tok_len - 1]
        return extn


    def get_file_category(self, url):
        if self.is_resume(url):
            return 'resume'
        else:
            return 'other'


    def get_file_tech_stack(self, text):
        tech_stacks = re.findall(regex_all.tech_stack, text.lower())
        t_s_dict = {}
        ts_list = []
        for t_s in tech_stacks:
            if t_s in t_s_dict.keys():
                t_s_dict[t_s] = t_s_dict[t_s] + 1
            else:
                t_s_dict[t_s] = 1
        for key, value in t_s_dict.items():
            ts_list.append("{}({})".format(key,value))
        return ', '.join(ts_list)



    def get_file_total_experience(self, text):
        exp = 0
        years = re.findall(regex_all.year_in_resume, text)
        # print(years)
        s_years = sorted(years)
        # print(s_years)
        y_len = len(s_years)
        if y_len > 1:
            exp = int(s_years[y_len-1]) - int(s_years[0])
        return exp

    def get_bold_italics_underlined_texts(self, url_or_path):
        b_i_u = {'bold': '',
                 'italics': '',
                 'underlined': ''}
        if url_or_path.endswith('.pdf'):
            word_path = extract_pdf2.convert_to_word(url_or_path)
            if word_path is not None:
                document = Document(word_path)
                bolds,italics,underlined = [],[],[]
                for para in document.paragraphs:
                    for run in para.runs:
                        if run.bold:
                            bolds.append(run.text)
                        if run.italic:
                            italics.append(run.text)
                        if run.font.underline:
                            underlined.append(run.text)
                os.remove(word_path)
                b_i_u['bold'] = ','.join(bolds)
                b_i_u['italics'] = ','.join(italics)
                b_i_u['underlined'] = ','.join(underlined)
        return b_i_u

    def get_all_capitalized_keywords(self, text):
        return ', '.join(re.findall(regex_all.capitalized_words, text))

    def get_all_keywords(self, text):
        return ', '.join(re.findall(regex_all.resume, text, re.IGNORECASE))

    def get_font_size_ratios(self, url_or_path):
        word_path = extract_pdf2.convert_to_word(url_or_path)
        doc = docx.Document(word_path)
        for p in doc.paragraphs:
            name = p.style.font.name
            size = p.style.font.size
            print(name, size)

    def get_no_of_pages(self, test_url):
        pass


domains = [
        ]

def main():
    # pdf_online = "https://coreyms.com/portfolio/docs/Corey-Schafer-Resume.pdf"
    pdfs_online =[
        "https://www.troupcountyga.gov/Content/Documents/employment/two-rivers-rcd-executive-directory-ad.pdf",
        "https://www.troupcountyga.gov/Content/Documents/finance/cafr-16.pdf",
        "https://www.troupcountyga.gov/Content/Documents/finance/cafr-15.pdf",
        "https://www.troupcountyga.gov/Content/Documents/finance/cafr-13.pdf"
    ]
    # rtf_online = "https://file-examples-com.github.io/uploads/2019/09/file-sample_100kB.rtf"
    # pptx_online = "https://scholar.harvard.edu/files/torman_personal/files/samplepptx.pptx"
    # docx_online= "https://file-examples-com.github.io/uploads/2017/02/file-sample_100kB.docx"
    # xlsx_online= "https://file-examples-com.github.io/uploads/2017/02/file_example_XLSX_10.xlsx"
    # xls_online= "https://file-examples-com.github.io/uploads/2017/02/file_example_XLS_10.xls"
    # ppt_online = "https://file-examples-com.github.io/uploads/2017/08/file_example_PPT_250kB.ppt"
    # doc_online = "https://file-examples-com.github.io/uploads/2017/02/file-sample_100kB.doc"
    #file available in downloads directory
    # zip_file = "https://transfer.sh/bn3eY0/resume.zip"
    # zip_file = "https://transfer.sh/vtLon5/zip_file_2.zip"
    resume_pattern = re.compile(regex_all.resume)
    m = Metadata(resume_pattern)
    # print(m.get_file_tech_stack('bro javA enabdajd JavA dude kafKa bro javA dude kafKa dudelinuxlinuxmanlinux'))
    # print(m.get_file_total_experience('5454545454545 2012 1985 dkjkldjg.2020.jkfdjkf'))
    # logging.basicConfig(
    #     format='%(asctime)s %(levelname)-8s %(message)s',
    #     level=logging.INFO,
    #     datefmt='%Y-%m-%d %H:%M:%S',
    #     filename='metadata.log')
    # m.req_resp_size_logger(m.get_resp_url("sony").resp)
    for pdf in pdfs_online:
        print("crawling:{}".format(pdf))
        text = m.get_text_from_pdf(pdf)
        matches = re.finditer(m.resume_pattern, text)
        for match in matches:
            print(match)
    # print(m.does_zip_contain_resume(zip_file))
    # print(m.get_text_from_zip_file_url(zip_file))
    # resp = requests.get("http://www.africau.edu/images/default/sample.pdf")
    # print(m.get_text_from_pdf(pdf_online))
    # if pdf_len > 0: print("pdf parsing success: len: {}".format(pdf_len))
    # pdf_len = len(m.get_text_from_pdf_new(pdf_online))
    # if pdf_len > 0: print("pdf parsing success: len: {}".format(pdf_len))
    # if len(m.get_text_from_othr_file(rtf_online)) > 0: print("rtf parsing success")
    # if len(m.get_text_from_pptx('/tmp/zip_file/samplepptx.pptx')) > 0: print("pptx parsing success")
    # if len(m.get_text_from_docx(docx_online)) > 0: print("docx parsing success")
    # if len(m.get_text_from_excel(xlsx_online)) > 0: print("xlsx parsing success")
    # if len(m.get_text_from_excel('/tmp/zip_file/file_example_XLSX_10.xlsx')) > 0: print("xlsx parsing success")
    # df = m.get_df_from_excel('/tmp/zip_file/file_example_XLSX_10.xlsx')
    # print (df)
    # rslt = ''
    # for i,j in df.iterrows():
    #     col_row = '{}:{}'.format(i,j)
    #     # print(col_row)
    #     if None != re.search('kiro',col_row,re.IGNORECASE):
    #         rslt = "found:"+col_row
    #         break
    # if '' == rslt:
    #     print("not found")
    # else:
    #     print (rslt)
        # df[i].astype('str').apply(lambda x: print('found:'+df[i].name) if re.search('a',x,re.IGNORECASE) else 'pass')
    # if len(m.get_text_from_xls(xls_online)) > 0: print("xls parsing success")
    # if len(m.get_text_from_ppt(ppt_online)) > 0: print("ppt parsing success")
    # if len(m.get_text_from_doc(doc_online)) > 0: print("doc parsing success")

    # text = m.get_text_from_pptx(pptx_online)
    # print(m.get_ip_address("qa.bigparser.com"))
    # print(m.get_resp_url("coreyms").url)
    # print(m.get_language(m.get_resp_url("coreyms").resp))
    # print(m.get_ip_country("coreyms.com"))


if __name__ == "__main__":
    main()
